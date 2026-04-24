# =========================
# Dokumentläsning
# =========================
#
# Läser in olika dokumentformat och gör om dem till text som kan användas i resten av pipelinen.


import csv
from pathlib import Path

import openpyxl
from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from src.config import DOCUMENT_LOADING_SETTINGS


# =========================
# TXT
# =========================
#
# Läser vanliga textfiler med den encoding som är angiven i config.

def load_txt(file_path: str) -> str:
    """
    Läser in text från en TXT-fil.
    """
    txt_settings = DOCUMENT_LOADING_SETTINGS.get("txt", {})
    encoding = txt_settings.get("encoding")

    if not encoding:
        raise ValueError("DOCUMENT_LOADING_SETTINGS saknar 'txt.encoding' i config.py.")

    with open(file_path, "r", encoding=encoding) as f:
        return f.read()


# =========================
# DOCX
# =========================
#
# Hämtar text från både paragrafer och tabeller.

def load_docx(file_path: str) -> str:
    """
    Läser in text från en DOCX-fil genom att extrahera paragrafer och tabeller.
    """
    doc = Document(file_path)
    texts = []

    # Vanliga paragrafer (löpande text, rubriker)
    for p in doc.paragraphs:
        if p.text.strip():
            texts.append(p.text)
            
    # Text i tabeller (vanligt i CV-mallar, rapporter)
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                texts.append(" | ".join(row_texts))  # Bevarar tabellstruktur läsbart

    return "\n\n".join(texts)


# =========================
# PDF
# =========================
#
# Läser textbaserade PDF:er och försöker rätta vanliga teckenkodningsfel.
# Inskannade PDF:er kräver OCR och stöds inte ännu.

def load_pdf(file_path: str) -> str:
    """
    Läser in text från en textbaserad PDF-fil.
    """
    pdf_settings = DOCUMENT_LOADING_SETTINGS.get("pdf", {})
    encodings_to_try = pdf_settings.get("encodings_to_try")

    if not encodings_to_try:
        raise ValueError(
            "DOCUMENT_LOADING_SETTINGS saknar 'pdf.encodings_to_try' i config.py."
        )

    reader = PdfReader(file_path)
    pages = []   

    for i, page in enumerate(reader.pages):
        text = page.extract_text(extraction_mode="layout")
        if not text or not text.strip():
            continue

        # Försök rätta kodningen med varje par i tur och ordning
        fixed_text = text
        for encode_as, decode_as in encodings_to_try:
            try:
                candidate = text.encode(encode_as).decode(decode_as)
                # Välj kandidaten om den innehåller fler giltiga svenska tecken
                if candidate.count("å") + candidate.count("ä") + candidate.count("ö") + \
                   candidate.count("Å") + candidate.count("Ä") + candidate.count("Ö") > \
                   fixed_text.count("å") + fixed_text.count("ä") + fixed_text.count("ö") + \
                   fixed_text.count("Å") + fixed_text.count("Ä") + fixed_text.count("Ö"):
                    fixed_text = candidate
            except (UnicodeDecodeError, UnicodeEncodeError):
                continue

        pages.append(f"[Sida {i + 1}]\n{fixed_text.strip()}")

    if not pages:
        raise ValueError(
            "Kunde inte extrahera text från PDF:en. Filen kan vara en inskannad bild — OCR stöds inte ännu."
        )

    return "\n\n".join(pages)


# =========================
# PPTX
# =========================
#
# Hämtar text från slides, tabeller och talaranteckningar.

def load_pptx(file_path: str) -> str:
    """
    Läser in text från en PPTX-fil genom att extrahera slide-text, tabeller och anteckningar.
    """
    prs = Presentation(file_path)
    texts = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = [f"[Slide {i}]"]

        for shape in slide.shapes:
            # Textrutor och rubriker
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text and para.text.strip():
                            slide_texts.append(para.text.strip())
            except Exception:
                pass

            # Tabeller i slides
            try:
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text and cell.text.strip()
                        ]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))
            except Exception:
                pass

        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text if notes_frame and notes_frame.text else ""
                notes_text = notes_text.strip()
                if notes_text:
                    slide_texts.append(f"[Anteckningar]: {notes_text}")
        except Exception:
            pass

        if len(slide_texts) > 1:
            texts.append("\n".join(slide_texts))

    return "\n\n".join(texts)


# =========================
# CSV-hjälpfunktioner
# =========================
#
# Hjälper CSV-läsningen när automatisk delimiter-detektering inte räcker.

def _get_fallback_csv_delimiter(sample: str, delimiters: str) -> str:
    """
    Väljer fallback-separator för CSV om csv.Sniffer inte kan avgöra dialect.

    Strategi:
    - räkna hur många gånger varje tillåten separator förekommer i sample
    - välj separatorn med högst förekomst
    - om ingen förekommer, använd komma som standard
    """
    if not isinstance(sample, str):
        raise TypeError("sample måste vara en sträng.")

    if not isinstance(delimiters, str):
        raise TypeError("delimiters måste vara en sträng.")

    if not delimiters:
        raise ValueError("delimiters kan inte vara tom.")

    delimiter_counts = {
        delimiter: sample.count(delimiter)
        for delimiter in delimiters
    }

    best_delimiter = max(delimiter_counts, key=delimiter_counts.get)

    if delimiter_counts[best_delimiter] == 0:
        return ","

    return best_delimiter


# =========================
# CSV
# =========================
#
# Läser in CSV-rader och gör om dem till läsbar text för sökning och RAG.

def load_csv(file_path: str) -> str:
    """
    Läser in text från en CSV-fil och omvandlar rader till läsbar text.
    """
    csv_settings = DOCUMENT_LOADING_SETTINGS.get("csv", {})
    encoding = csv_settings.get("encoding")
    sample_size = csv_settings.get("sample_size")
    delimiters = csv_settings.get("sniffer_delimiters")

    if not encoding:
        raise ValueError("DOCUMENT_LOADING_SETTINGS saknar 'csv.encoding' i config.py.")
    
    if not sample_size:
        raise ValueError("DOCUMENT_LOADING_SETTINGS saknar 'csv.sample_size' i config.py.")
    
    if not delimiters:
        raise ValueError(
            "DOCUMENT_LOADING_SETTINGS saknar 'csv.sniffer_delimiters' i config.py."
        )

    texts = []

    with open(file_path, "r", encoding=encoding) as f:
        # Detekterar automatiskt separator (, eller ;)
        sample = f.read(sample_size)
        f.seek(0)

        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=delimiters)
        except csv.Error:
            fallback_delimiter = _get_fallback_csv_delimiter(sample, delimiters)

            class SimpleDialect(csv.Dialect):
                delimiter = fallback_delimiter
                quotechar = '"'
                doublequote = True
                skipinitialspace = False
                lineterminator = "\n"
                quoting = csv.QUOTE_MINIMAL

            dialect = SimpleDialect

        reader = csv.DictReader(f, dialect=dialect)

        headers = reader.fieldnames or []
        texts.append("Kolumner: " + ", ".join(headers))

        for row in reader:
            # Varje rad omvandlas till läsbar text: "Kolumn: Värde"
            row_text = " | ".join(
                f"{k}: {v.strip()}" for k, v in row.items() if v and v.strip()
            )
            if row_text:
                texts.append(row_text)

    return "\n\n".join(texts)


# =========================
# XLSX-laddning
# =========================
#
# Denna del innehåller logik för att läsa in Excel-filer genom att extrahera
# innehåll från alla flikar och göra raderna läsbara som text.

def load_xlsx(file_path: str) -> str:
    """
    Läser in text från en XLSX-fil genom att extrahera innehåll från alla flikar.
    """
    wb = openpyxl.load_workbook(file_path, data_only=True)
    texts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_texts = [f"[Flik: {sheet_name}]"]

        for row in ws.iter_rows(values_only=True):
            # Filtrera tomma rader
            row_values = [
                str(cell).strip() 
                for cell in row 
                if cell is not None and str(cell).strip()
            ]
            if row_values:
                sheet_texts.append(" | ".join(row_values))

        if len(sheet_texts) > 1:
            texts.append("\n\n".join(sheet_texts))

    return "\n\n".join(texts)


# =========================
# Routing till rätt dokumentladdare
# =========================
#
# Denna del ansvarar för att välja rätt loader utifrån filändelse och skicka
# filen vidare till korrekt funktion för textutvinning.

def load_document(file_path: str) -> str:
    """
    Väljer rätt loader baserat på filändelse och returnerar extraherad text.
    """
    if not isinstance(file_path, (str, Path)):
        raise TypeError("file_path måste vara en sträng eller Path.")
    
    if not str(file_path).strip():
        raise ValueError("file_path saknas.")

    suffix = Path(file_path).suffix.lower()

    loaders = {
        ".txt": load_txt,
        ".docx": load_docx,
        ".pdf": load_pdf,
        ".pptx": load_pptx,
        ".csv": load_csv,
        ".xlsx": load_xlsx,
    }

    if not suffix:
        raise ValueError("Filen saknar filändelse och kan därför inte matchas mot en loader.")

    if suffix not in loaders:
        raise ValueError(f"Filtypen stöds inte ännu: {suffix}")

    return loaders[suffix](file_path)
